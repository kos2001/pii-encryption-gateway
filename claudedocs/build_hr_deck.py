#!/usr/bin/env python3
"""Build an HR-facing explainer deck for the PII Encryption Gateway skill.

Visual style: Samsung DS-division brand *approximation* (not an official
template file — none was available). Samsung Blue primary, clean executive
16:9 layout, Korean-friendly typography. Re-run to regenerate the .pptx.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- Samsung DS brand approximation ----
BLUE       = RGBColor(0x14, 0x28, 0xA0)   # Samsung Blue (primary)
BLUE_DARK  = RGBColor(0x0B, 0x14, 0x5A)   # title band
ACCENT     = RGBColor(0x3E, 0x5C, 0xD0)   # accent / links
SKY        = RGBColor(0xE8, 0xED, 0xFB)   # light card / chip
CARD       = RGBColor(0xF4, 0xF6, 0xFB)   # neutral card
TEXT       = RGBColor(0x22, 0x22, 0x22)
MUTED      = RGBColor(0x6B, 0x72, 0x80)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GREEN      = RGBColor(0x1A, 0x9E, 0x5E)
AMBER      = RGBColor(0xE5, 0x97, 0x00)
RED        = RGBColor(0xD6, 0x45, 0x41)

HEAD_FONT = "맑은 고딕"   # Malgun Gothic — Korean corporate standard
BODY_FONT = "맑은 고딕"
MONO_FONT = "Consolas"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def _set_font(run, name):
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)


def text(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         wrap=True, line_spacing=1.0):
    """runs: list of lines; each line is str OR list of (txt,size,color,bold,font) tuples."""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(2)
    for i, line in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        segs = [line] if isinstance(line, tuple) else (line if isinstance(line, list) else [line])
        if isinstance(line, str):
            segs = [(line, 16, TEXT, False, BODY_FONT)]
        for seg in segs:
            txt, size, color, bold, font = seg
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            _set_font(r, font)
    return tb


def rect(slide, l, t, w, h, fill, line=None, shape=MSO_SHAPE.RECTANGLE, line_w=1.0):
    sp = slide.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def shape_text(sp, runs, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE):
    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(6)
    tf.margin_top = tf.margin_bottom = Pt(4)
    for i, line in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        for seg in ([line] if isinstance(line, tuple) else line):
            txt, size, color, bold, font = seg
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
            _set_font(r, font)


def header(slide, title, kicker=None):
    rect(slide, 0, 0, 13.333, 1.15, WHITE)
    rect(slide, 0, 1.15, 13.333, 0.045, BLUE)              # accent rule
    if kicker:
        text(slide, 0.6, 0.18, 11, 0.3, [[(kicker, 12, ACCENT, True, HEAD_FONT)]])
    text(slide, 0.6, 0.42, 12, 0.7, [[(title, 26, BLUE_DARK, True, HEAD_FONT)]],
         anchor=MSO_ANCHOR.MIDDLE)
    footer(slide)


def footer(slide):
    text(slide, 0.6, 7.06, 7, 0.3,
         [[("대외비 · Confidential", 9, MUTED, False, BODY_FONT)]])
    text(slide, 6.0, 7.06, 6.7, 0.3,
         [[("SAMSUNG ", 9, BLUE, True, HEAD_FONT), ("DS  |  인사팀 대상 설명자료", 9, MUTED, False, BODY_FONT)]],
         align=PP_ALIGN.RIGHT)


def chip(slide, l, t, w, label, fill=SKY, fg=BLUE):
    sp = rect(slide, l, t, w, 0.42, fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    shape_text(sp, [[(label, 12, fg, True, BODY_FONT)]])
    return sp


# ============================== SLIDE 1: COVER ==============================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.333, 7.5, BLUE_DARK)
rect(s, 0, 0, 0.28, 7.5, ACCENT)                # left accent bar
text(s, 0.9, 0.7, 6, 0.4, [[("SAMSUNG ", 16, WHITE, True, HEAD_FONT),
                            ("DS", 16, RGBColor(0x9F,0xB3,0xF2), True, HEAD_FONT)]])
text(s, 0.9, 2.5, 11.5, 2.0, [
    [("개인정보 보호 게이트웨이", 44, WHITE, True, HEAD_FONT)],
    [("PII Encryption Gateway", 22, RGBColor(0x9F,0xB3,0xF2), False, HEAD_FONT)],
], line_spacing=1.1)
text(s, 0.95, 4.7, 11, 0.6,
     [[("인사 데이터를 AI에 ", 18, RGBColor(0xD7,0xDE,0xF7), False, BODY_FONT),
       ("그대로 노출하지 않고", 18, WHITE, True, BODY_FONT),
       (" 안전하게 활용하는 법", 18, RGBColor(0xD7,0xDE,0xF7), False, BODY_FONT)]])
text(s, 0.95, 6.4, 11, 0.4,
     [[("인사팀 대상 설명자료  ·  2026.06  ·  대외비", 13, RGBColor(0x9F,0xB3,0xF2), False, BODY_FONT)]])

# ===================== SLIDE 2: WHY (the problem) =====================
s = prs.slides.add_slide(BLANK)
header(s, "왜 필요한가요?", "배경")
text(s, 0.6, 1.5, 12.1, 0.6,
     [[("인사 데이터에는 ", 18, TEXT, False, BODY_FONT),
       ("가장 민감한 개인정보", 18, RED, True, BODY_FONT),
       ("가 들어 있습니다.", 18, TEXT, False, BODY_FONT)]])
for i, (t_, d) in enumerate([
        ("연봉·급여", "급여대장"), ("주민등록번호", "신원 정보"),
        ("계좌번호", "급여 이체"), ("전화·이메일", "연락처"), ("근태", "지각·결근")]):
    x = 0.6 + i * 2.42
    c = rect(s, x, 2.25, 2.25, 1.0, SKY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    shape_text(c, [[(t_, 15, BLUE, True, BODY_FONT)], [(d, 11, MUTED, False, BODY_FONT)]])
text(s, 0.6, 3.75, 12.1, 0.6,
     [[("그런데 이 데이터로 ", 18, TEXT, False, BODY_FONT),
       ("AI에게 안내문·통지·리포트", 18, BLUE, True, BODY_FONT),
       ("를 맡기려면?", 18, TEXT, False, BODY_FONT)]])
prob = rect(s, 0.6, 4.5, 12.13, 1.9, RGBColor(0xFD,0xF1,0xF0),
            line=RED, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
shape_text(prob, [
    [("⚠  위험: 원본을 그대로 AI에 넣으면 민감정보가 모델에 노출됩니다", 18, RED, True, BODY_FONT)],
    [("최민준 · 770324-1809570 · 5,350만원 · 신한 617-1434-688508 …", 14, MUTED, False, MONO_FONT)],
    [("한 번 노출된 개인정보는 되돌릴 수 없고, 유출 시 법적·평판 리스크가 큽니다.", 14, TEXT, False, BODY_FONT)],
], align=PP_ALIGN.LEFT)

# ===================== SLIDE 3: WHAT IT IS (one-liner + analogy) =====================
s = prs.slides.add_slide(BLANK)
header(s, "한 장 요약: 무엇인가요?", "개념")
big = rect(s, 0.6, 1.5, 12.13, 1.35, BLUE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
shape_text(big, [
    [("민감한 값을 ‘대리 번호표(토큰)’로 바꿔 AI에 보여주고,", 19, WHITE, True, BODY_FONT)],
    [("작업이 끝나면 담당자 키로만 원래 값을 되돌려 주는 안전 통로입니다.", 19, WHITE, True, BODY_FONT)],
])
text(s, 0.6, 3.15, 12, 0.4, [[("🔐  은행 금고에 비유하면", 15, ACCENT, True, BODY_FONT)]])
for i, (h, d) in enumerate([
    ("금고에 보관", "원본 민감정보는 암호화해 ‘금고(vault)’에 넣습니다"),
    ("번호표로 작업", "AI는 실제 값 대신 번호표(토큰)만 보고 일합니다"),
    ("열쇠 가진 사람만", "담당자 키가 있어야 원본을 꺼내볼 수 있습니다")]):
    x = 0.6 + i * 4.04
    c = rect(s, x, 3.6, 3.85, 2.5, CARD, line=RGBColor(0xE2,0xE6,0xF0),
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    num = rect(s, x + 0.25, 3.85, 0.6, 0.6, BLUE, shape=MSO_SHAPE.OVAL)
    shape_text(num, [[(str(i + 1), 18, WHITE, True, HEAD_FONT)]])
    text(s, x + 0.25, 4.65, 3.4, 0.5, [[(h, 17, BLUE_DARK, True, BODY_FONT)]])
    text(s, x + 0.25, 5.15, 3.4, 1.0, [[(d, 13, TEXT, False, BODY_FONT)]], line_spacing=1.1)

# ===================== SLIDE 4: HOW IT WORKS (3-step flow) =====================
s = prs.slides.add_slide(BLANK)
header(s, "어떻게 동작하나요? — 3단계", "동작 원리")
steps = [
    ("1. 보호 (Protect)", "원본 → 토큰\n원본은 키로 암호화",
     "5,350만원  →  [[SALARY:3f9a2c1d]]"),
    ("2. 작업 (Work)", "AI는 토큰만 보고\n안내문·통지·리포트 작성",
     "“[[NAME]]님의 연봉은 [[SALARY]]…”"),
    ("3. 복원 (Reveal)", "담당자 키로 토큰 →\n원본 (인가자만)",
     "[[SALARY:3f9a2c1d]]  →  5,350만원"),
]
for i, (h, d, ex) in enumerate(steps):
    x = 0.6 + i * 4.35
    c = rect(s, x, 1.7, 3.95, 3.7, WHITE, line=BLUE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, line_w=1.5)
    rect(s, x, 1.7, 3.95, 0.7, BLUE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, x, 1.78, 3.95, 0.55, [[(h, 16, WHITE, True, BODY_FONT)]], align=PP_ALIGN.CENTER)
    text(s, x + 0.3, 2.65, 3.35, 1.3, [[(ln, 14, TEXT, False, BODY_FONT)] for ln in d.split("\n")],
         align=PP_ALIGN.CENTER, line_spacing=1.15)
    exb = rect(s, x + 0.25, 4.25, 3.45, 0.95, SKY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    shape_text(exb, [[(seg, 11, BLUE_DARK, False, MONO_FONT)] for seg in [ex]])
    if i < 2:
        ar = rect(s, x + 3.98, 3.25, 0.34, 0.6, ACCENT, shape=MSO_SHAPE.CHEVRON)
banner = rect(s, 0.6, 5.75, 12.13, 0.85, SKY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
shape_text(banner, [[("핵심: 작업 중 AI 화면(중간본)에는 ", 15, TEXT, False, BODY_FONT),
                     ("실제 값이 단 한 번도 등장하지 않습니다.", 15, BLUE, True, BODY_FONT)]])

# ===================== SLIDE 5: WHAT IT PROTECTS =====================
s = prs.slides.add_slide(BLANK)
header(s, "무엇을, 어디까지 보호하나요?", "보호 범위")
text(s, 0.6, 1.45, 12, 0.4, [[("자동으로 찾아 가리는 민감정보", 16, BLUE_DARK, True, BODY_FONT)]])
ents = ["주민등록번호", "연봉·급여", "계좌번호", "전화번호", "이메일", "카드번호",
        "사업자등록번호", "IP 주소", "직원 이름(명부)"]
for i, e in enumerate(ents):
    chip(s, 0.6 + (i % 5) * 2.45, 1.95 + (i // 5) * 0.55, 2.3, e)
text(s, 0.6, 3.25, 12, 0.4, [[("네 가지 방식으로 빠짐없이 탐지", 16, BLUE_DARK, True, BODY_FONT)]])
ways = [
    ("① 컬럼 이름", "‘연봉’, ‘주민번호’ 같은 열을 인식"),
    ("② 컬럼 형태", "이름이 달라도 값 모양으로 자동 판별"),
    ("③ 값 모양", "문장 속에 섞인 번호·연락처도 탐지"),
    ("④ 명부 이름", "담당자 명부의 직원 이름을 정확히 가림"),
]
for i, (h, d) in enumerate(ways):
    x = 0.6 + (i % 2) * 6.15
    y = 3.75 + (i // 2) * 1.05
    c = rect(s, x, y, 5.95, 0.92, CARD, line=RGBColor(0xE2,0xE6,0xF0), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, x + 0.25, y + 0.1, 2.4, 0.7, [[(h, 14, BLUE, True, BODY_FONT)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, x + 2.5, y + 0.1, 3.3, 0.7, [[(d, 12.5, TEXT, False, BODY_FONT)]], anchor=MSO_ANCHOR.MIDDLE)
note = rect(s, 0.6, 5.95, 12.13, 0.65, SKY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
shape_text(note, [[("엑셀·CSV 명부뿐 아니라 ", 14, TEXT, False, BODY_FONT),
                   ("메모·이메일 초안 같은 문서(.txt/.md)", 14, BLUE, True, BODY_FONT),
                   ("도 그대로 보호됩니다.", 14, TEXT, False, BODY_FONT)]])

# ===================== SLIDE 6: BENEFITS FOR HR =====================
s = prs.slides.add_slide(BLANK)
header(s, "인사팀에게 무엇이 좋은가요?", "기대 효과")
benefits = [
    ("🛡", "안전", "민감정보가 AI에 0건 노출.\n유출 리스크를 원천 차단"),
    ("⚡", "간편", "평소처럼 안내문·통지문 작성.\n보호는 도구가 알아서 처리"),
    ("🔑", "되돌림", "담당자 키로만 원복.\n키 없으면 복원 불가 = 접근통제"),
    ("💸", "효율", "대규모에선 토큰도 더 적게 사용\n(원본 전체를 읽지 않으므로)"),
]
for i, (ic, h, d) in enumerate(benefits):
    x = 0.6 + i * 3.06
    c = rect(s, x, 1.7, 2.9, 2.7, WHITE, line=RGBColor(0xE2,0xE6,0xF0), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    icb = rect(s, x + 1.05, 1.95, 0.8, 0.8, SKY, shape=MSO_SHAPE.OVAL)
    shape_text(icb, [[(ic, 24, BLUE, False, BODY_FONT)]])
    text(s, x, 2.85, 2.9, 0.4, [[(h, 18, BLUE_DARK, True, BODY_FONT)]], align=PP_ALIGN.CENTER)
    text(s, x + 0.2, 3.3, 2.5, 1.0, [[(ln, 12.5, TEXT, False, BODY_FONT)] for ln in d.split("\n")],
         align=PP_ALIGN.CENTER, line_spacing=1.1)
# evidence band
ev = rect(s, 0.6, 4.85, 12.13, 1.65, BLUE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
shape_text(ev, [
    [("검증 결과 (동일 과제 비교)", 14, RGBColor(0xC8,0xD3,0xF5), True, BODY_FONT)],
    [("게이트웨이 사용 100%  vs  미사용 62.5%", 22, WHITE, True, BODY_FONT)],
    [("개인 안내·통지 작업에서 미사용은 매번 실제 PII가 노출 — 게이트웨이는 누출 0 유지", 13, RGBColor(0xD7,0xDE,0xF7), False, BODY_FONT)],
])

# ===================== SLIDE 7: EXAMPLE (HR scenario) =====================
s = prs.slides.add_slide(BLANK)
header(s, "실제로 이렇게 쓰입니다", "사용 예시")
text(s, 0.6, 1.45, 12, 0.5,
     [[("예) 사번 E0007 직원에게 ", 16, TEXT, False, BODY_FONT),
       ("2026년 연봉 안내 메일", 16, BLUE, True, BODY_FONT), (" 작성", 16, TEXT, False, BODY_FONT)]])
# draft (tokens) card
d1 = rect(s, 0.6, 2.05, 5.95, 3.6, CARD, line=RGBColor(0xE2,0xE6,0xF0), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
rect(s, 0.6, 2.05, 5.95, 0.6, MUTED, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, 0.6, 2.12, 5.95, 0.45, [[("AI가 보는 중간본 (토큰)", 14, WHITE, True, BODY_FONT)]], align=PP_ALIGN.CENTER)
text(s, 0.9, 2.85, 5.4, 2.7, [
    [("받는사람: [[NAME:369de402]]", 12.5, TEXT, False, MONO_FONT)],
    [("", 6, TEXT, False, MONO_FONT)],
    [("안녕하세요, [[NAME:369de402]]님.", 12.5, TEXT, False, MONO_FONT)],
    [("2026년 확정 연봉은", 12.5, TEXT, False, MONO_FONT)],
    [("[[SALARY:28314826]] 입니다.", 12.5, TEXT, False, MONO_FONT)],
    [("", 6, TEXT, False, MONO_FONT)],
    [("→ 실제 이름·연봉 없음 ✓", 12.5, GREEN, True, BODY_FONT)],
], line_spacing=1.15)
arrow = rect(s, 6.62, 3.55, 0.55, 0.6, ACCENT, shape=MSO_SHAPE.CHEVRON)
shape_text(arrow, [[("키", 12, WHITE, True, BODY_FONT)]])
# final (real) card
d2 = rect(s, 7.3, 2.05, 5.43, 3.6, WHITE, line=BLUE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, line_w=1.5)
rect(s, 7.3, 2.05, 5.43, 0.6, BLUE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, 7.3, 2.12, 5.43, 0.45, [[("담당자만 보는 최종본 (원본 복원)", 14, WHITE, True, BODY_FONT)]], align=PP_ALIGN.CENTER)
text(s, 7.6, 2.85, 4.9, 2.7, [
    [("받는사람: 신다은", 12.5, TEXT, False, BODY_FONT)],
    [("", 6, TEXT, False, BODY_FONT)],
    [("안녕하세요, 신다은님.", 12.5, TEXT, False, BODY_FONT)],
    [("2026년 확정 연봉은", 12.5, TEXT, False, BODY_FONT)],
    [("3,194만원 입니다.", 12.5, TEXT, True, BODY_FONT)],
    [("", 6, TEXT, False, BODY_FONT)],
    [("→ 인가된 담당자에게만 ✓", 12.5, BLUE, True, BODY_FONT)],
], line_spacing=1.15)
text(s, 0.6, 5.85, 12, 0.5,
     [[("같은 방식으로 ", 14, TEXT, False, BODY_FONT),
       ("근태 통지·상여금 입금 안내·부서 리포트", 14, BLUE, True, BODY_FONT),
       (" 등 일상 업무에 그대로 적용됩니다.", 14, TEXT, False, BODY_FONT)]])

# ===================== SLIDE 8: LIMITS (honest) =====================
s = prs.slides.add_slide(BLANK)
header(s, "꼭 알아둘 점 (솔직한 한계)", "유의사항")
items = [
    ("토큰으로는 계산 불가", "‘평균 연봉’ 같은 집계는 토큰으로 못 합니다. 원본 대상 스크립트로 계산해 결과 숫자만 활용합니다."),
    ("같은 값 = 같은 번호표", "그룹핑이 가능한 대신, 누가 같은 값을 갖는지(동일성)는 드러납니다. 값 자체는 키 없이 복원 불가."),
    ("문장 속 이름은 명부 필요", "패턴이 없는 ‘이름’은 담당자 명부를 줘야 가립니다. 명부에 없는 외부 인물 이름은 탐지 못 함."),
    ("키 관리가 핵심", "복원은 담당자 키로만 가능 → 키를 안전하게 보관·관리해야 합니다. 키 분실 시 복원 불가."),
]
for i, (h, d) in enumerate(items):
    y = 1.6 + i * 1.32
    c = rect(s, 0.6, y, 12.13, 1.15, CARD, line=RGBColor(0xE2,0xE6,0xF0), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    bar = rect(s, 0.6, y, 0.12, 1.15, AMBER)
    text(s, 0.95, y + 0.12, 4.0, 0.95, [[(h, 15, BLUE_DARK, True, BODY_FONT)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 5.0, y + 0.12, 7.5, 0.95, [[(d, 12.5, TEXT, False, BODY_FONT)]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)

# ===================== SLIDE 9: GETTING STARTED / CLOSE =====================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.333, 7.5, BLUE_DARK)
rect(s, 0, 0, 0.28, 7.5, ACCENT)
text(s, 0.9, 0.8, 11, 0.5, [[("시작하기", 30, WHITE, True, HEAD_FONT)]])
steps2 = [
    ("1  보호", "민감 파일을 게이트웨이에 통과 → 토큰본 + 암호화 금고 생성"),
    ("2  작업", "토큰본으로 평소처럼 안내문·통지·리포트 작성"),
    ("3  복원", "담당자 키로 최종본만 원복하여 발송"),
]
for i, (h, d) in enumerate(steps2):
    y = 1.9 + i * 1.0
    text(s, 1.0, y, 2.2, 0.6, [[(h, 18, RGBColor(0x9F,0xB3,0xF2), True, HEAD_FONT)]])
    text(s, 3.1, y + 0.03, 9.0, 0.6, [[(d, 15, WHITE, False, BODY_FONT)]])
box = rect(s, 0.95, 5.2, 11.4, 1.25, RGBColor(0x16, 0x22, 0x6E),
           line=ACCENT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
shape_text(box, [
    [("문의 · 도입 지원", 13, RGBColor(0x9F,0xB3,0xF2), True, BODY_FONT)],
    [("담당자 키 발급과 적용 방법은 데이터보호 담당 부서로 문의해 주세요.", 15, WHITE, False, BODY_FONT)],
], align=PP_ALIGN.LEFT)
text(s, 0.95, 6.85, 11.4, 0.4,
     [[("※ 본 자료의 디자인은 Samsung DS 브랜드를 근사한 것이며, 데이터는 모두 합성(가상) 예시입니다.",
        10, RGBColor(0x8A,0x97,0xC9), False, BODY_FONT)]])

import os
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "PII_게이트웨이_인사팀_설명.pptx")
prs.save(out)
print("saved:", out, "| slides:", len(list(prs.slides)))
